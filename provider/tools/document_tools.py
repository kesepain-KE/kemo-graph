"""本地文档到 Markdown 的安全、高保真转换与导入工具。

转换器只处理可直接提取文本的文档。扫描版 PDF/OCR 不在本模块职责内，
遇到没有文本层的 PDF 时会返回明确错误，避免生成看似成功的空文档。
"""

from __future__ import annotations

import codecs
import csv
import hashlib
import json
import os
import re
import unicodedata
import xml.dom.minidom
import xml.etree.ElementTree as ET
from collections import Counter
from io import StringIO
from pathlib import Path
from typing import Callable, Iterable, Iterator
from uuid import uuid4


class DocumentConversionError(RuntimeError):
    """文档无法读取或转换。"""


def convert_pdf(path: Path | str) -> str:
    """提取文本型 PDF，清理重复页眉页脚并保留页边界。"""

    source = _validate_input_file(path, {".pdf"})
    try:
        import pdfplumber

        with pdfplumber.open(source) as document:
            raw_pages = [
                (page.extract_text(x_tolerance=2, y_tolerance=3, layout=False) or "")
                for page in document.pages
            ]
    except Exception as exc:
        raise DocumentConversionError(f"无法转换 PDF：{source}: {exc}") from exc

    if not any(page.strip() for page in raw_pages):
        raise DocumentConversionError(
            f"PDF 没有可提取的文本层：{source}。扫描版 PDF 需要由主智能体先执行 OCR。"
        )

    repeated_edges = _repeated_pdf_edge_lines(raw_pages)
    pages: list[str] = []
    for index, raw_page in enumerate(raw_pages, start=1):
        cleaned = _clean_pdf_page(raw_page, repeated_edges)
        if cleaned:
            pages.append(f"## 第 {index} 页\n\n{cleaned}")
    return _finalize_markdown("\n\n".join(pages), source)


def convert_docx(path: Path | str) -> str:
    """按正文顺序转换 DOCX 的标题、列表、段落、链接与表格。"""

    source = _validate_input_file(path, {".docx"})
    try:
        from docx import Document

        document = Document(source)
    except Exception as exc:
        raise DocumentConversionError(f"无法读取 DOCX：{source}: {exc}") from exc

    blocks: list[str] = []
    for block in _iter_docx_blocks(document):
        if _is_docx_paragraph(block):
            rendered = _render_docx_paragraph(block)
        else:
            rows = [
                [_escape_markdown_cell(cell.text) for cell in row.cells]
                for row in block.rows
            ]
            rendered = _markdown_table(rows) if rows else ""
        if rendered:
            blocks.append(rendered)
    return _finalize_markdown("\n\n".join(blocks), source)


def convert_html(path: Path | str) -> str:
    """清理网页噪声后将 HTML 主体转换为 Markdown。"""

    source = _validate_input_file(path, {".html", ".htm"})
    html = _read_text(source)
    try:
        from bs4 import BeautifulSoup
        from markdownify import markdownify

        soup = BeautifulSoup(html, "html.parser")
        for node in soup.select(
            "script, style, noscript, template, svg, canvas, iframe, form, nav, button, input"
        ):
            node.decompose()
        main = soup.find("article") or soup.find("main") or soup.body or soup
        markdown = markdownify(
            str(main),
            heading_style="ATX",
            bullets="-",
        )
    except Exception as exc:
        raise DocumentConversionError(f"无法转换 HTML：{source}: {exc}") from exc
    return _finalize_markdown(markdown, source)


def convert_txt(path: Path | str) -> str:
    """读取普通文本并保留段落结构，不再将全文包装为代码块。"""

    source = _validate_input_file(path, {".txt", ".log"})
    return _finalize_markdown(_read_text(source), source)


def convert_rst(path: Path | str) -> str:
    """通过 docutils 将 RST 转 HTML，再转 Markdown。"""

    source = _validate_input_file(path, {".rst"})
    rst = _read_text(source)
    try:
        from docutils.core import publish_parts
        from markdownify import markdownify

        html = publish_parts(source=rst, writer_name="html5")["html_body"]
        markdown = markdownify(html, heading_style="ATX")
    except ImportError as exc:
        raise DocumentConversionError(
            "转换 RST 需要 docutils，请执行：python -m pip install -r requirements.txt"
        ) from exc
    except Exception as exc:
        raise DocumentConversionError(f"无法转换 RST：{source}: {exc}") from exc
    return _finalize_markdown(markdown, source)


def convert_csv(path: Path | str) -> str:
    """自动识别编码与分隔符，将 CSV/TSV 转换为 Markdown 表格。"""

    source = _validate_input_file(path, {".csv", ".tsv"})
    text = _read_text(source)
    delimiter = "\t" if source.suffix.casefold() == ".tsv" else None
    try:
        sample = text[:65536]
        dialect = None
        if delimiter is None:
            try:
                dialect = csv.Sniffer().sniff(sample, delimiters=",\t;|")
            except csv.Error:
                delimiter = ","
        reader = (
            csv.reader(StringIO(text), dialect=dialect)
            if dialect
            else csv.reader(StringIO(text), delimiter=delimiter)
        )
        rows = [list(row) for row in reader]
    except csv.Error as exc:
        raise DocumentConversionError(f"无法解析表格文本：{source}: {exc}") from exc
    return _finalize_markdown(_markdown_table(rows), source)


def convert_spreadsheet(path: Path | str) -> str:
    """将 XLSX/XLSM/XLS 各工作表按原始顺序转换为 Markdown 表格。"""

    source = _validate_input_file(path, {".xlsx", ".xlsm", ".xls"})
    try:
        if source.suffix.casefold() == ".xls":
            sheets = _read_legacy_xls(source)
        else:
            sheets = _read_openpyxl_workbook(source)
    except Exception as exc:
        raise DocumentConversionError(f"无法转换电子表格：{source}: {exc}") from exc

    blocks: list[str] = []
    for sheet_name, rows in sheets:
        blocks.append(f"## {_escape_heading(sheet_name)}")
        blocks.append(_markdown_table(rows) if rows else "_空工作表_ ")
    return _finalize_markdown("\n\n".join(blocks), source)


def convert_pptx(path: Path | str) -> str:
    """按幻灯片和页面内位置提取 PPTX 文本及表格。"""

    source = _validate_input_file(path, {".pptx"})
    try:
        from pptx import Presentation

        presentation = Presentation(source)
    except Exception as exc:
        raise DocumentConversionError(f"无法读取 PPTX：{source}: {exc}") from exc

    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        blocks = [f"## 幻灯片 {index}"]
        shapes = sorted(
            slide.shapes,
            key=lambda item: (int(item.top), int(item.left)),
        )
        for shape in shapes:
            if getattr(shape, "has_table", False):
                rows = [
                    [_escape_markdown_cell(cell.text) for cell in row.cells]
                    for row in shape.table.rows
                ]
                if rows:
                    blocks.append(_markdown_table(rows))
                continue
            text = str(getattr(shape, "text", "") or "").strip()
            if text:
                blocks.append(text)
        slides.append("\n\n".join(blocks))
    return _finalize_markdown("\n\n---\n\n".join(slides), source)


def convert_epub(path: Path | str) -> str:
    """按书脊顺序提取 EPUB 章节并转换为 Markdown。"""

    source = _validate_input_file(path, {".epub"})
    try:
        from bs4 import BeautifulSoup
        from ebooklib import epub
        from markdownify import markdownify

        book = epub.read_epub(str(source), options={"ignore_ncx": True})
        chapters: list[str] = []
        for chapter_id, _linear in book.spine:
            item = book.get_item_with_id(chapter_id)
            if item is None:
                continue
            soup = BeautifulSoup(item.get_content(), "html.parser")
            for node in soup.select("script, style, nav, svg"):
                node.decompose()
            body = soup.body or soup
            rendered = markdownify(str(body), heading_style="ATX").strip()
            if rendered:
                chapters.append(rendered)
    except Exception as exc:
        raise DocumentConversionError(f"无法转换 EPUB：{source}: {exc}") from exc
    return _finalize_markdown("\n\n---\n\n".join(chapters), source)


def convert_rtf(path: Path | str) -> str:
    """提取 RTF 文本并保留自然段。"""

    source = _validate_input_file(path, {".rtf"})
    try:
        from striprtf.striprtf import rtf_to_text

        markdown = rtf_to_text(_read_text(source), errors="ignore")
    except Exception as exc:
        raise DocumentConversionError(f"无法转换 RTF：{source}: {exc}") from exc
    return _finalize_markdown(markdown, source)


def convert_data(path: Path | str) -> str:
    """规范化 JSON/JSONL/YAML/XML 数据文档并用对应代码围栏呈现。"""

    source = _validate_input_file(
        path, {".json", ".jsonl", ".ndjson", ".yaml", ".yml", ".xml"}
    )
    text = _read_text(source)
    suffix = source.suffix.casefold()
    try:
        if suffix == ".json":
            rendered = json.dumps(json.loads(text), ensure_ascii=False, indent=2)
            language = "json"
        elif suffix in {".jsonl", ".ndjson"}:
            objects = [json.loads(line) for line in text.splitlines() if line.strip()]
            rendered = "\n".join(json.dumps(item, ensure_ascii=False) for item in objects)
            language = "jsonl"
        elif suffix in {".yaml", ".yml"}:
            import yaml

            rendered = yaml.safe_dump(
                yaml.safe_load(text), allow_unicode=True, sort_keys=False, default_flow_style=False
            )
            language = "yaml"
        else:
            parsed = ET.fromstring(text)
            rendered = xml.dom.minidom.parseString(
                ET.tostring(parsed, encoding="utf-8")
            ).toprettyxml(indent="  ", encoding=None)
            language = "xml"
    except Exception as exc:
        raise DocumentConversionError(f"无法解析结构化数据：{source}: {exc}") from exc
    return _finalize_markdown(f"```{language}\n{rendered.strip()}\n```", source)


def import_document(
    path: Path | str,
    external_dir: Path | str,
    *,
    destination_name: str | None = None,
) -> dict[str, object]:
    """检测格式、转换并原子写入 external/markdown。"""

    source = _validate_input_file(path, set(SUPPORTED_DOCUMENT_SUFFIXES))
    markdown_dir = _validate_output_dir(external_dir)
    suffix = source.suffix.lower()
    if suffix in {".md", ".markdown"}:
        markdown = _finalize_markdown(_read_text(source), source)
        format_name = "markdown"
    else:
        converter = _CONVERTERS.get(suffix)
        if converter is None:
            raise DocumentConversionError(f"不支持的文档格式：{suffix or '<无扩展名>'}")
        markdown = converter(source)
        format_name = suffix.removeprefix(".")

    markdown_dir.mkdir(parents=True, exist_ok=True)
    safe_stem = _safe_filename(source.stem)
    destination = (
        _validated_destination(markdown_dir, destination_name)
        if destination_name is not None
        else markdown_dir / f"{safe_stem}.md"
    )
    encoded = markdown.encode("utf-8")
    if destination_name is None and destination.exists() and destination.read_bytes() != encoded:
        digest = hashlib.sha256(str(source).encode("utf-8") + b"\0" + encoded).hexdigest()[:10]
        destination = markdown_dir / f"{safe_stem}-{digest}.md"

    if not destination.exists() or destination.read_bytes() != encoded:
        destination.parent.mkdir(parents=True, exist_ok=True)
        temporary = destination.with_name(f".{destination.name}.{uuid4().hex}.tmp")
        try:
            temporary.write_bytes(encoded)
            os.replace(temporary, destination)
        except OSError as exc:
            try:
                temporary.unlink(missing_ok=True)
            except OSError:
                pass
            raise DocumentConversionError(f"无法写入 Markdown：{destination}: {exc}") from exc

    return {
        "source_path": str(destination.resolve()),
        "markdown_relative_path": destination.relative_to(markdown_dir).as_posix(),
        "format": format_name,
        "ok": True,
    }


def _validate_input_file(value: Path | str, allowed_suffixes: set[str]) -> Path:
    path = _raw_absolute_path(value, "path")
    if path.suffix.lower() not in allowed_suffixes:
        expected = "、".join(sorted(allowed_suffixes))
        raise DocumentConversionError(f"文件格式不匹配，期望 {expected}：{path}")
    resolved = path.resolve(strict=False)
    if not resolved.exists():
        raise FileNotFoundError(f"文件不存在：{resolved}")
    if not resolved.is_file():
        raise DocumentConversionError(f"路径不是普通文件：{resolved}")
    try:
        with resolved.open("rb") as stream:
            stream.read(1)
    except OSError as exc:
        raise DocumentConversionError(f"文件不可读：{resolved}: {exc}") from exc
    return resolved


def _validate_output_dir(value: Path | str) -> Path:
    root = _raw_absolute_path(value, "external_dir").resolve(strict=False)
    return root if root.name.casefold() == "markdown" else root / "markdown"


def _validated_destination(root: Path, value: str) -> Path:
    if not isinstance(value, str) or not value.strip():
        raise DocumentConversionError("destination_name 必须是非空相对路径")
    relative = Path(value.strip())
    if relative.is_absolute() or ".." in relative.parts:
        raise DocumentConversionError("destination_name 必须是安全的相对路径")
    if relative.suffix.casefold() != ".md":
        raise DocumentConversionError("destination_name 必须以 .md 结尾")
    destination = (root / relative).resolve()
    try:
        destination.relative_to(root.resolve())
    except ValueError as exc:
        raise DocumentConversionError("destination_name 超出 Markdown 目录") from exc
    return destination


def _raw_absolute_path(value: Path | str, name: str) -> Path:
    if not isinstance(value, (str, Path)):
        raise TypeError(f"{name} 必须是路径字符串")
    raw = str(value).strip()
    if not raw:
        raise ValueError(f"{name} 不能为空")
    path = Path(raw).expanduser()
    if ".." in path.parts:
        raise DocumentConversionError(f"拒绝包含 '..' 的路径：{raw}")
    if not path.is_absolute():
        raise DocumentConversionError(f"必须使用绝对路径：{raw}")
    return path


def _read_text(path: Path) -> str:
    """解码常见中英文文本编码，并拒绝明显的二进制/乱码结果。"""

    try:
        data = path.read_bytes()
    except OSError as exc:
        raise DocumentConversionError(f"无法读取文件：{path}: {exc}") from exc
    if not data:
        return ""

    bom_candidates = (
        (codecs.BOM_UTF8, "utf-8-sig"),
        (codecs.BOM_UTF32_LE, "utf-32"),
        (codecs.BOM_UTF32_BE, "utf-32"),
        (codecs.BOM_UTF16_LE, "utf-16"),
        (codecs.BOM_UTF16_BE, "utf-16"),
    )
    for bom, encoding in bom_candidates:
        if data.startswith(bom):
            try:
                return data.decode(encoding)
            except UnicodeError:
                break

    try:
        return data.decode("utf-8")
    except UnicodeError:
        pass

    candidates: list[tuple[float, str, str]] = []
    try:
        from charset_normalizer import from_bytes

        best = from_bytes(data).best()
        if best is not None and best.encoding:
            guessed_encoding = best.encoding.casefold().replace("-", "_")
            # UTF-16/32 无 BOM 时很容易把短 GB18030 文本误判成成对宽字符。
            if not guessed_encoding.startswith(("utf_16", "utf_32")):
                decoded = str(best)
                candidates.append((_decoded_text_penalty(decoded), best.encoding, decoded))
    except Exception:
        pass

    for encoding in ("gb18030", "big5", "shift_jis", "cp1252"):
        try:
            decoded = data.decode(encoding)
        except UnicodeError:
            continue
        candidates.append((_decoded_text_penalty(decoded), encoding, decoded))
    if not candidates:
        raise DocumentConversionError(f"无法识别文本编码：{path}")
    penalty, encoding, decoded = min(candidates, key=lambda item: item[0])
    if penalty >= 0.35:
        raise DocumentConversionError(f"文件疑似二进制或编码已损坏：{path}（候选编码 {encoding}）")
    return decoded


def _decoded_text_penalty(text: str) -> float:
    if not text:
        return 0.0
    controls = sum(
        1 for char in text if unicodedata.category(char) == "Cc" and char not in "\n\r\t"
    )
    replacement = text.count("\ufffd")
    mojibake = sum(text.count(marker) for marker in ("锟斤拷", "ï»¿", "Ã", "Â"))
    return (controls * 5 + replacement * 10 + mojibake * 3) / max(len(text), 1)


def _finalize_markdown(markdown: str, source: Path) -> str:
    text = unicodedata.normalize("NFC", str(markdown))
    text = text.replace("\r\n", "\n").replace("\r", "\n").replace("\x00", "")
    text = "".join(
        char
        for char in text
        if unicodedata.category(char) != "Cc" or char in "\n\t"
    )
    lines = [line.rstrip() for line in text.split("\n")]
    text = "\n".join(lines)
    text = re.sub(r"\n{4,}", "\n\n\n", text).strip()
    if not text:
        raise DocumentConversionError(f"转换结果为空：{source}")
    return text + "\n"


def _repeated_pdf_edge_lines(pages: list[str]) -> set[str]:
    if len(pages) < 3:
        return set()
    candidates: Counter[str] = Counter()
    for page in pages:
        lines = [line.strip() for line in page.splitlines() if line.strip()]
        for line in set(lines[:2] + lines[-2:]):
            if 2 <= len(line) <= 160:
                candidates[line] += 1
    threshold = max(3, (len(pages) + 1) // 2)
    return {line for line, count in candidates.items() if count >= threshold}


def _clean_pdf_page(page: str, repeated_edges: set[str]) -> str:
    lines = [line.strip() for line in page.replace("\u00a0", " ").splitlines()]
    lines = [line for line in lines if line and line not in repeated_edges]
    text = "\n".join(lines)
    text = re.sub(r"(?<=[A-Za-z])-[ \t]*\n(?=[a-z])", "", text)
    return text.strip()


def _iter_docx_blocks(document: object) -> Iterator[object]:
    from docx.document import Document as DocumentType
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P

    if not isinstance(document, DocumentType):
        return
    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, document)
        elif isinstance(child, CT_Tbl):
            yield Table(child, document)


def _is_docx_paragraph(block: object) -> bool:
    from docx.text.paragraph import Paragraph

    return isinstance(block, Paragraph)


def _render_docx_paragraph(paragraph: object) -> str:
    from docx.oxml.ns import qn

    parts: list[str] = []
    for child in paragraph._p.iterchildren():
        if child.tag == qn("w:r"):
            parts.append(_render_docx_run_element(child, paragraph))
        elif child.tag == qn("w:hyperlink"):
            label = "".join(child.itertext()).strip()
            relationship_id = child.get(qn("r:id"))
            target = ""
            if relationship_id:
                relationship = paragraph.part.rels.get(relationship_id)
                target = str(getattr(relationship, "target_ref", "") or "")
            parts.append(f"[{label}]({target})" if label and target else label)
    text = "".join(parts).strip() or paragraph.text.strip()
    if not text:
        return ""
    style_name = str(getattr(paragraph.style, "name", "") or "")
    heading = re.fullmatch(r"(?:Heading|标题)\s*([1-6])", style_name, flags=re.IGNORECASE)
    if heading:
        return f"{'#' * int(heading.group(1))} {text}"
    if paragraph._p.pPr is not None and paragraph._p.pPr.numPr is not None:
        marker = "1." if "number" in style_name.casefold() else "-"
        return f"{marker} {text}"
    return text


def _render_docx_run_element(element: object, paragraph: object) -> str:
    from docx.text.run import Run

    run = Run(element, paragraph)
    text = run.text
    if not text:
        return ""
    if run.bold:
        text = f"**{text}**"
    if run.italic:
        text = f"*{text}*"
    if run.font.strike:
        text = f"~~{text}~~"
    return text


def _read_openpyxl_workbook(source: Path) -> list[tuple[str, list[list[str]]]]:
    from openpyxl import load_workbook

    workbook = load_workbook(source, read_only=True, data_only=False, keep_links=False)
    try:
        return [
            (
                sheet.title,
                [
                    _trim_trailing_empty(
                        [_spreadsheet_value(cell.value) for cell in row]
                    )
                    for row in sheet.iter_rows()
                ],
            )
            for sheet in workbook.worksheets
        ]
    finally:
        workbook.close()


def _read_legacy_xls(source: Path) -> list[tuple[str, list[list[str]]]]:
    import xlrd

    workbook = xlrd.open_workbook(str(source), on_demand=True)
    try:
        return [
            (
                sheet.name,
                [
                    _trim_trailing_empty(
                        [
                            _spreadsheet_value(sheet.cell_value(row, col))
                            for col in range(sheet.ncols)
                        ]
                    )
                    for row in range(sheet.nrows)
                ],
            )
            for sheet in workbook.sheets()
        ]
    finally:
        workbook.release_resources()


def _spreadsheet_value(value: object) -> str:
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value)


def _trim_trailing_empty(row: list[str]) -> list[str]:
    while row and not row[-1].strip():
        row.pop()
    return row


def _markdown_table(rows: Iterable[Iterable[object]]) -> str:
    normalized_rows = [
        [_escape_markdown_cell(str(value)) for value in row]
        for row in rows
    ]
    normalized_rows = [row for row in normalized_rows if any(cell for cell in row)]
    if not normalized_rows:
        return ""
    width = max(len(row) for row in normalized_rows)
    header = normalized_rows[0] + [""] * (width - len(normalized_rows[0]))
    body = [row + [""] * (width - len(row)) for row in normalized_rows[1:]]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in range(width)) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


def _escape_markdown_cell(value: str) -> str:
    return (
        str(value)
        .replace("\\", "\\\\")
        .replace("|", "\\|")
        .replace("\r", "")
        .replace("\n", "<br>")
        .strip()
    )


def _escape_heading(value: str) -> str:
    return str(value).replace("\n", " ").strip() or "未命名工作表"


def _safe_filename(value: str) -> str:
    cleaned = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "_", value).strip(" .")
    return cleaned or "document"


SUPPORTED_DOCUMENT_SUFFIXES = frozenset(
    {
        ".pdf",
        ".docx",
        ".pptx",
        ".xlsx",
        ".xlsm",
        ".xls",
        ".html",
        ".htm",
        ".epub",
        ".rtf",
        ".txt",
        ".log",
        ".rst",
        ".csv",
        ".tsv",
        ".json",
        ".jsonl",
        ".ndjson",
        ".yaml",
        ".yml",
        ".xml",
        ".md",
        ".markdown",
    }
)

_SUPPORTED_SUFFIXES = set(SUPPORTED_DOCUMENT_SUFFIXES)

_CONVERTERS: dict[str, Callable[[Path | str], str]] = {
    ".pdf": convert_pdf,
    ".docx": convert_docx,
    ".pptx": convert_pptx,
    ".xlsx": convert_spreadsheet,
    ".xlsm": convert_spreadsheet,
    ".xls": convert_spreadsheet,
    ".html": convert_html,
    ".htm": convert_html,
    ".epub": convert_epub,
    ".rtf": convert_rtf,
    ".txt": convert_txt,
    ".log": convert_txt,
    ".rst": convert_rst,
    ".csv": convert_csv,
    ".tsv": convert_csv,
    ".json": convert_data,
    ".jsonl": convert_data,
    ".ndjson": convert_data,
    ".yaml": convert_data,
    ".yml": convert_data,
    ".xml": convert_data,
}


__all__ = [
    "DocumentConversionError",
    "SUPPORTED_DOCUMENT_SUFFIXES",
    "convert_pdf",
    "convert_docx",
    "convert_html",
    "convert_txt",
    "convert_rst",
    "convert_csv",
    "convert_spreadsheet",
    "convert_pptx",
    "convert_epub",
    "convert_rtf",
    "convert_data",
    "import_document",
]
