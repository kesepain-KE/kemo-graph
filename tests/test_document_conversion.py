"""多编码与多媒体文档转换回归测试。"""

from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from provider.tools.document_tools import (
    DocumentConversionError,
    convert_csv,
    convert_data,
    convert_docx,
    convert_epub,
    convert_html,
    convert_pptx,
    convert_rtf,
    convert_spreadsheet,
    convert_txt,
)


class DocumentConversionTests(unittest.TestCase):
    def test_common_chinese_encodings_are_decoded_without_mojibake(self) -> None:
        with tempfile.TemporaryDirectory() as temporary_dir:
            root = Path(temporary_dir)
            samples = {
                "gb18030.txt": ("中文段落：知识图谱与向量检索。", "gb18030"),
                "big5.txt": ("繁體中文：知識圖譜與向量檢索。", "big5"),
                "utf16.txt": ("UTF-16 文本：正常显示。", "utf-16"),
            }
            for filename, (expected, encoding) in samples.items():
                path = root / filename
                path.write_bytes(expected.encode(encoding))
                self.assertIn(expected, convert_txt(path))

    def test_csv_detects_chinese_encoding_and_semicolon_dialect(self) -> None:
        with tempfile.TemporaryDirectory() as temporary_dir:
            source = Path(temporary_dir) / "名单.csv"
            source.write_bytes("姓名;城市\n小明;上海\n".encode("gb18030"))
            converted = convert_csv(source)
            self.assertIn("| 姓名 | 城市 |", converted)
            self.assertIn("| 小明 | 上海 |", converted)

    def test_docx_preserves_paragraph_and_table_order(self) -> None:
        from docx import Document

        with tempfile.TemporaryDirectory() as temporary_dir:
            source = Path(temporary_dir) / "ordered.docx"
            document = Document()
            document.add_heading("文档标题", level=1)
            document.add_paragraph("表格之前")
            table = document.add_table(rows=2, cols=2)
            table.rows[0].cells[0].text = "项目"
            table.rows[0].cells[1].text = "内容"
            table.rows[1].cells[0].text = "图谱"
            table.rows[1].cells[1].text = "关系"
            document.add_paragraph("表格之后")
            document.save(source)

            converted = convert_docx(source)
            self.assertIn("# 文档标题", converted)
            self.assertLess(converted.index("表格之前"), converted.index("| 项目 | 内容 |"))
            self.assertLess(converted.index("| 项目 | 内容 |"), converted.index("表格之后"))

    def test_html_removes_script_and_navigation_noise(self) -> None:
        with tempfile.TemporaryDirectory() as temporary_dir:
            source = Path(temporary_dir) / "page.html"
            source.write_bytes(
                "<html><body><nav>菜单噪声</nav><main><h1>正文</h1><p>有效内容</p>"
                "<script>danger()</script></main></body></html>".encode("gb18030")
            )
            converted = convert_html(source)
            self.assertIn("# 正文", converted)
            self.assertIn("有效内容", converted)
            self.assertNotIn("菜单噪声", converted)
            self.assertNotIn("danger", converted)

    def test_spreadsheet_and_presentation_keep_visible_content(self) -> None:
        from openpyxl import Workbook
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as temporary_dir:
            root = Path(temporary_dir)
            workbook_path = root / "data.xlsx"
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "知识表"
            sheet.append(["节点", "权重"])
            sheet.append(["阅读", 3])
            workbook.save(workbook_path)
            spreadsheet = convert_spreadsheet(workbook_path)
            self.assertIn("## 知识表", spreadsheet)
            self.assertIn("| 节点 | 权重 |", spreadsheet)
            self.assertIn("| 阅读 | 3 |", spreadsheet)

            presentation_path = root / "slides.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "知识图谱"
            slide.placeholders[1].text = "节点与关系"
            presentation.save(presentation_path)
            converted_pptx = convert_pptx(presentation_path)
            self.assertIn("知识图谱", converted_pptx)
            self.assertIn("节点与关系", converted_pptx)

    def test_structured_data_and_rtf_are_normalized(self) -> None:
        with tempfile.TemporaryDirectory() as temporary_dir:
            root = Path(temporary_dir)
            json_path = root / "data.json"
            json_path.write_text('{"主题":"阅读","权重":2}', encoding="utf-8")
            self.assertIn('"主题": "阅读"', convert_data(json_path))

            yaml_path = root / "data.yaml"
            yaml_path.write_text("主题: 阅读\n权重: 2\n", encoding="utf-8")
            self.assertIn("主题: 阅读", convert_data(yaml_path))

            rtf_path = root / "note.rtf"
            rtf_path.write_text(r"{\rtf1\ansi Knowledge Graph\par Vector Search}", encoding="ascii")
            converted_rtf = convert_rtf(rtf_path)
            self.assertIn("Knowledge Graph", converted_rtf)
            self.assertIn("Vector Search", converted_rtf)

    def test_epub_uses_spine_order_and_keeps_chapter_text(self) -> None:
        from ebooklib import epub

        with tempfile.TemporaryDirectory() as temporary_dir:
            source = Path(temporary_dir) / "book.epub"
            book = epub.EpubBook()
            book.set_identifier("conversion-test")
            book.set_title("测试书籍")
            book.set_language("zh")
            chapter = epub.EpubHtml(
                title="第一章", file_name="chapter.xhtml", lang="zh"
            )
            chapter.content = (
                "<html><body><h1>第一章</h1><p>知识图谱正文。</p></body></html>"
            )
            book.add_item(chapter)
            book.add_item(epub.EpubNcx())
            book.add_item(epub.EpubNav())
            book.spine = ["nav", chapter]
            epub.write_epub(str(source), book)

            converted = convert_epub(source)
            self.assertIn("# 第一章", converted)
            self.assertIn("知识图谱正文", converted)

    def test_empty_document_is_rejected_instead_of_silent_success(self) -> None:
        with tempfile.TemporaryDirectory() as temporary_dir:
            source = Path(temporary_dir) / "empty.txt"
            source.write_bytes(b"")
            with self.assertRaisesRegex(DocumentConversionError, "转换结果为空"):
                convert_txt(source)


if __name__ == "__main__":
    unittest.main()
