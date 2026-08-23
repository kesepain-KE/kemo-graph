"""轻量 MarkItDown 核心入口与真实文档回归测试。"""

from __future__ import annotations

import subprocess
import sys
import tempfile
import unittest
from io import BytesIO
from pathlib import Path

from markitdown import (
    DocumentConverter,
    DocumentConverterResult,
    MarkItDown,
    UnsafeInputError,
)
from markitdown._stream_info import StreamInfo


class MarkItDownCoreTests(unittest.TestCase):
    def test_python_api_and_local_only_boundary(self) -> None:
        with tempfile.TemporaryDirectory() as temporary_dir:
            source = Path(temporary_dir) / "note.txt"
            source.write_text("第一行\n\n第二行", encoding="utf-8")
            result = MarkItDown().convert(source)
            self.assertEqual(result.markdown, result.text_content)
            self.assertIn("第一行", result.text_content)
            self.assertEqual(result.converter, "PlainTextConverter")

            with self.assertRaises(UnsafeInputError):
                MarkItDown().convert("https://example.com/note.txt")

    def test_real_pdf_docx_xlsx_and_pptx_files(self) -> None:
        with tempfile.TemporaryDirectory() as temporary_dir:
            root = Path(temporary_dir)

            pdf_path = root / "sample.pdf"
            pdf_path.write_bytes(_minimal_pdf("PDF knowledge graph"))
            pdf_result = MarkItDown().convert(pdf_path)
            self.assertIn("PDF knowledge graph", pdf_result.text_content)

            from docx import Document

            docx_path = root / "sample.docx"
            document = Document()
            document.add_heading("Document heading", level=1)
            document.add_paragraph("Document body")
            document.save(docx_path)
            docx_result = MarkItDown().convert(docx_path)
            self.assertIn("Document heading", docx_result.text_content)
            self.assertIn("Document body", docx_result.text_content)

            from openpyxl import Workbook

            xlsx_path = root / "sample.xlsx"
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "Data"
            sheet.append(["Node", "Weight"])
            sheet.append(["Graph", 3])
            workbook.save(xlsx_path)
            xlsx_result = MarkItDown().convert(xlsx_path)
            self.assertIn("## Data", xlsx_result.text_content)
            self.assertIn("| Node | Weight |", xlsx_result.text_content)

            from pptx import Presentation

            pptx_path = root / "sample.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Presentation heading"
            slide.placeholders[1].text = "Presentation body"
            presentation.save(pptx_path)
            pptx_result = MarkItDown().convert(pptx_path)
            self.assertIn("Presentation heading", pptx_result.text_content)
            self.assertIn("Presentation body", pptx_result.text_content)

    def test_converter_registry_is_extensible(self) -> None:
        class CustomConverter(DocumentConverter):
            extensions = frozenset({".custom"})
            priority = 100

            def convert(self, source: Path, info: StreamInfo) -> DocumentConverterResult:
                return DocumentConverterResult(
                    title="custom",
                    text_content="# Custom\n",
                    converter="CustomConverter",
                )

        with tempfile.TemporaryDirectory() as temporary_dir:
            source = Path(temporary_dir) / "sample.custom"
            source.write_text("payload", encoding="utf-8")
            converter = MarkItDown(converters=[CustomConverter()])
            result = converter.convert(source)
            self.assertEqual(result.converter, "CustomConverter")
            self.assertEqual(result.text_content, "# Custom\n")

    def test_binary_stream_uses_same_dispatcher(self) -> None:
        result = MarkItDown().convert_stream(
            BytesIO(b"stream content"),
            file_extension="txt",
            filename="stream.txt",
        )
        self.assertIn("stream content", result.text_content)
        self.assertEqual(result.source_path, "stream.txt")

    def test_eml_is_a_local_document_format(self) -> None:
        with tempfile.TemporaryDirectory() as temporary_dir:
            source = Path(temporary_dir) / "message.eml"
            source.write_text(
                "Subject: Knowledge note\n"
                "From: sender@example.com\n"
                "To: receiver@example.com\n"
                "\n"
                "The graph body.",
                encoding="utf-8",
            )
            result = MarkItDown().convert(source)
            self.assertIn("Knowledge note", result.text_content)
            self.assertIn("The graph body", result.text_content)

    def test_module_cli_writes_markdown(self) -> None:
        with tempfile.TemporaryDirectory() as temporary_dir:
            root = Path(temporary_dir)
            source = root / "cli.txt"
            output = root / "cli.md"
            source.write_text("CLI conversion", encoding="utf-8")
            completed = subprocess.run(
                [sys.executable, "-m", "markitdown", str(source), "-o", str(output)],
                cwd=Path(__file__).parents[1],
                capture_output=True,
                text=True,
                check=False,
            )
            self.assertEqual(completed.returncode, 0, completed.stderr)
            self.assertIn("CLI conversion", output.read_text(encoding="utf-8"))


def _minimal_pdf(text: str) -> bytes:
    stream = f"BT /F1 18 Tf 72 720 Td ({text}) Tj ET".encode("ascii")
    bodies = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
        b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    pdf = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, body in enumerate(bodies, start=1):
        offsets.append(len(pdf))
        pdf.extend(f"{index} 0 obj\n".encode("ascii"))
        pdf.extend(body)
        pdf.extend(b"\nendobj\n")
    xref_offset = len(pdf)
    pdf.extend(f"xref\n0 {len(bodies) + 1}\n".encode("ascii"))
    pdf.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        pdf.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    pdf.extend(
        f"trailer\n<< /Size {len(bodies) + 1} /Root 1 0 R >>\nstartxref\n{xref_offset}\n%%EOF\n".encode(
            "ascii"
        )
    )
    return bytes(pdf)


if __name__ == "__main__":
    unittest.main()
