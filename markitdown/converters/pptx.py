"""PowerPoint 文本和表格 Converter。"""

from __future__ import annotations

from pathlib import Path

from .._base_converter import DocumentConverter, DocumentConverterResult
from .._exceptions import DocumentConversionError, MissingOptionalDependencyError
from .._stream_info import StreamInfo
from .._utils import finalize_markdown, markdown_table, title_from_path


class PptxConverter(DocumentConverter):
    extensions = frozenset({".pptx"})

    def convert(self, source: Path, info: StreamInfo) -> DocumentConverterResult:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise MissingOptionalDependencyError("python-pptx", "PowerPoint") from exc

        try:
            presentation = Presentation(source)
            slides: list[str] = []
            for index, slide in enumerate(presentation.slides, start=1):
                blocks = [f"## 幻灯片 {index}"]
                shapes = sorted(
                    slide.shapes,
                    key=lambda item: (int(item.top), int(item.left)),
                )
                for shape in shapes:
                    if getattr(shape, "has_table", False):
                        rows = [
                            [str(cell.text or "") for cell in row.cells]
                            for row in shape.table.rows
                        ]
                        if rows:
                            blocks.append(markdown_table(rows))
                        continue
                    text = str(getattr(shape, "text", "") or "").strip()
                    if text:
                        blocks.append(text)
                slides.append("\n\n".join(blocks))
            content = finalize_markdown("\n\n---\n\n".join(slides), source)
        except Exception as exc:
            raise DocumentConversionError(
                f"无法读取 PowerPoint：{source}: {exc}"
            ) from exc
        return DocumentConverterResult(
            title=title_from_path(source),
            text_content=content,
            source_path=str(source),
            converter=self.__class__.__name__,
        )


__all__ = ["PptxConverter"]
